import pandas as pd
import datetime as dt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
ACCENT = RGBColor(0xE8, 0x4C, 0x3D)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    return p

df = pd.read_csv('cv_hunt_career/CRM_marketer/dataset_crm.csv', on_bad_lines='skip')
if df.columns[0].startswith('Unnamed'):
    df = df.iloc[:, 1:]

df['Дата покупки'] = pd.to_datetime(df['Дата покупки'], format='%d.%m.%Y %H:%M:%S', errors='coerce')
df = df.dropna(subset=['Дата покупки', 'Контакт', 'Сумма'])

current_date = df['Дата покупки'].max() + dt.timedelta(days=1)

rfm = df.groupby('Контакт').agg({
    'Дата покупки': lambda x: (current_date - x.max()).days,
    'Карта': 'count',
    'Сумма': 'sum'
}).rename(columns={'Дата покупки': 'Recency', 'Карта': 'Frequency', 'Сумма': 'Monetary'})

rfm['Avg_Check'] = rfm['Monetary'] / rfm['Frequency']

total_clients = len(rfm)

q75_m = rfm['Monetary'].quantile(0.75)
seg1 = rfm[(rfm['Recency'] > 180) & (rfm['Frequency'] >= 2) & (rfm['Monetary'] >= q75_m)]

q75_avg = rfm['Avg_Check'].quantile(0.75)
seg2 = rfm[(rfm['Frequency'] == 1) & (rfm['Recency'] <= 180) & (rfm['Monetary'] >= q75_avg)]

q25_m = rfm['Monetary'].quantile(0.25)
seg3 = rfm[(rfm['Frequency'] <= 2) & (rfm['Recency'] > 90) & (rfm['Recency'] <= 365) & (rfm['Monetary'] >= q25_m) & (rfm['Monetary'] < q75_m)]

segments = {'Отток VIP': seg1, 'Новички с потенциалом': seg2, 'Спящие середняки': seg3}

deliver_rate = 0.45
cost_per_msg = 120
conv_rates = {'Отток VIP': 0.08, 'Новички с потенциалом': 0.12, 'Спящие середняки': 0.05}

results = []
for name, data in segments.items():
    base_size = len(data)
    delivered = int(base_size * deliver_rate)
    cost = delivered * cost_per_msg
    conv = conv_rates[name]
    checks = int(delivered * conv)
    avg_check = int(data['Avg_Check'].mean()) if len(data) > 0 else 0
    revenue = checks * avg_check
    roi = int((revenue - cost) / cost * 100) if cost > 0 else 0
    
    results.append({
        'name': name, 'base': base_size, 'delivered': delivered,
        'cost': cost, 'checks': checks, 'revenue': revenue,
        'roi': roi, 'avg_check': avg_check, 'conv': conv
    })

    print(f"[{name}] База: {base_size} чел. | Доставлено: {delivered} | Чеки: {checks} | Выручка: {revenue} | Затраты: {cost}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Титульник
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)
add_text(slide, 0.8, 1.8, 11, 1.2, "Стратегия CRM-коммуникаций", font_size=44, bold=True, color=WHITE)
add_text(slide, 0.8, 3.2, 11, 0.8, f"Анализ тестового датасета ({total_clients} клиентов)", font_size=22, bold=False, color=LIGHT_GRAY)

# Сегменты
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)
add_text(slide, 0.8, 0.8, 11, 0.8, "Сегменты (реальные цифры из таблицы)", font_size=32, bold=True, color=WHITE)
add_text(slide, 0.8, 2.0, 11, 0.5, f"1. Отток VIP: {results[0]['base']} человек. (Покупали 2+ раза, много тратили, не были больше 180 дней)", font_size=16)
add_text(slide, 0.8, 3.0, 11, 0.5, f"2. Новички с потенциалом: {results[1]['base']} человек. (1 покупка, большой чек, были недавно)", font_size=16)
add_text(slide, 0.8, 4.0, 11, 0.5, f"3. Спящие середняки: {results[2]['base']} человек. (1-2 покупки, средний чек, не были 3-12 мес)", font_size=16)

# Экономика
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)
add_text(slide, 0.8, 0.8, 11, 0.8, "Экономика по сегментам (WhatsApp 120₸)", font_size=32, bold=True, color=WHITE)
y_start = 2.0
headers = ["Сегмент", "Людей", "Доставлено(45%)", "Чеки", "Ср. чек", "Затраты", "Выручка", "ROI"]
col_widths = [2.5, 1.0, 1.5, 1.0, 1.5, 1.5, 1.5, 1.0]
x_start = 0.5
for i, header in enumerate(headers):
    x = x_start + sum(col_widths[:i])
    add_text(slide, x, y_start, col_widths[i], 0.35, header, font_size=13, bold=True, color=ACCENT)

for row_idx, r in enumerate(results):
    y = y_start + 0.5 + row_idx * 0.6
    vals = [r['name'], str(r['base']), str(r['delivered']), str(r['checks']), f"{r['avg_check']} ₸", f"{r['cost']} ₸", f"{r['revenue']} ₸", f"{r['roi']}%"]
    for i, val in enumerate(vals):
        x = x_start + sum(col_widths[:i])
        add_text(slide, x, y, col_widths[i], 0.35, val, font_size=13)

# Тексты
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACK)
add_text(slide, 0.8, 0.8, 11, 0.8, "Текст рассылки (WhatsApp 2-в-1)", font_size=32, bold=True, color=WHITE)
add_text(slide, 0.8, 2.0, 11, 1.5, "«INTERTOP-ты таңдағаныңызға рақмет! Жаңа аяқ киім топтамасына 5 000 бонус сыйлаймыз. Оларды ай соңына дейін қосымшада немесе сайтта жұмсап үлгеріңіз! 🎁\n\nСпасибо, что выбираете INTERTOP! Дарим вам 5 000 бонусов на покупку новой коллекции. Успейте потратить их до конца месяца в приложении или на сайте! 🎁»", font_size=14, color=LIGHT_GRAY)
add_text(slide, 0.8, 4.5, 11, 0.5, "Кнопка (CTA): Смотреть новинки / Жаңа топтаманы көру", font_size=16, bold=True, color=ACCENT)

prs.save('cv_hunt_career/CRM_marketer/Real_Numbers_Presentation.pptx')
print("\nУСПЕШНО! Файл Real_Numbers_Presentation.pptx создан.")
