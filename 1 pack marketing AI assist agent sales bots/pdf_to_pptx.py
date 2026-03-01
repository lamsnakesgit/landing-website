
import fitz
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_pptx(pdf_path, pptx_path):
    print(f"Converting {pdf_path} to {pptx_path}...")
    doc = fitz.open(pdf_path)
    prs = Presentation()
    
    # Set slide dimensions to match PDF (assuming all pages same size)
    # PPTX default is 10x7.5 inches. 
    # Let's get PDF first page size in points. 1 inch = 72 points.
    page0 = doc[0]
    width_pt = page0.rect.width
    height_pt = page0.rect.height
    
    # Update presentation slide size
    prs.slide_width = int(width_pt / 72 * 914400) # 914400 EMUs per inch
    prs.slide_height = int(height_pt / 72 * 914400)
    
    print(f"Slide Size: {width_pt} x {height_pt} points")

    for i, page in enumerate(doc):
        # Extract image
        pix = page.get_pixmap(dpi=150) # 150 dpi is decent for screen
        img_filename = f"temp_slide_{i}.png"
        pix.save(img_filename)
        
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add image filling the slide
        slide.shapes.add_picture(img_filename, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Clean up
        os.remove(img_filename)
        print(f"Added slide {i+1}/{len(doc)}")
        
    prs.save(pptx_path)
    print(f"Saved PowerPoint to {pptx_path}")

if __name__ == "__main__":
    pdf_to_pptx("Интеллектуальный_отдел_продаж_Branded.pdf", "Интеллектуальный_отдел_продаж_Editable.pptx")
